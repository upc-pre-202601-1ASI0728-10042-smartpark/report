# -*- coding: utf-8 -*-
"""Genera el Keynote TB2 y el deck About-the-Team (+ guiones) para Apex Twin / SmartPark."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = r'C:\codeProjects\SMARTPARK\report'
IMG = os.path.join(BASE, 'assets', 'images')

# --- Design system SmartPark ---
TEAL  = RGBColor(0x19,0x47,0x56)
TEAL2 = RGBColor(0x0B,0x7A,0x85)
CREAM = RGBColor(0xF2,0xEA,0xBC)
INK   = RGBColor(0x16,0x1A,0x1C)
CORAL = RGBColor(0xFF,0x3B,0x58)
WHITE = RGBColor(0xFF,0xFF,0xFF)
MUTED = RGBColor(0x5B,0x6B,0x70)
LIGHT = RGBColor(0xEF,0xF3,0xF4)

SW, SH = Inches(13.333), Inches(7.5)

MEMBERS = {
 'Elmer':  ('Riva Rodríguez, Elmer Augusto',   'Team Leader · Backend & Cloud'),
 'Camila': ('Sánchez Ríos, Camila Cristina',   'Landing & UI/UX'),
 'Britney':('Qqueso Rodríguez, Britney Delhy',  'App Móvil & Validación'),
 'Abel':   ('Valle Zuta, Abel Andrés',          'Frontend del Operador'),
 'Hernán': ('Morales Calderón, Hernan Emilio',  'Gemelo Digital & IoT'),
}

def ip(*parts):
    return os.path.join(IMG, *parts)

def img_size_in(path, bw, bh):
    """Devuelve (w,h) en pulgadas para encajar la imagen en (bw,bh) manteniendo aspecto."""
    iw, ih = Image.open(path).size
    a = iw / ih
    if bw / bh > a:
        h = bh; w = bh * a
    else:
        w = bw; h = bw / a
    return w, h

def rect(slide, shape_type, x, y, w, h, fill=None, line=None, line_w=None):
    sp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w or 1)
    sp.shadow.inherit = False
    return sp

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf

def setrun(p, text, size, color, bold=False, font='Calibri'):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = font
    return r

def add_image_fit(slide, path, bx, by, bw, bh):
    if not path or not os.path.exists(path):
        return
    w, h = img_size_in(path, bw, bh)
    x = bx + (bw - w) / 2; y = by + (bh - h) / 2
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

def video_box(slide):
    # Recuadro inferior derecho para el video del expositor
    bx, by, bw, bh = 10.15, 5.30, 2.95, 1.86
    box = rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh, fill=WHITE, line=TEAL, line_w=1.5)
    # franja superior tipo cámara
    rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, 0.42, fill=TEAL)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    setrun(p, 'Recuadro para el video', 11, MUTED, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    setrun(p2, 'del expositor', 11, MUTED)
    # etiqueta sobre la franja
    lab = textbox(slide, bx, by+0.04, bw, 0.36, MSO_ANCHOR.MIDDLE)
    lp = lab.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
    setrun(lp, 'CÁMARA EN VIVO', 9, WHITE, bold=True)

def footer(slide, presenter, idx, total, deck):
    fn, role = MEMBERS.get(presenter, (presenter, ''))
    tf = textbox(slide, 0.55, 7.02, 8.5, 0.4, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    setrun(p, f'{fn}', 11, TEAL, bold=True)
    setrun(p, f'   ·   {role}', 11, MUTED)
    # numero de slide
    tf2 = textbox(slide, 9.0, 7.02, 1.0, 0.4, MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
    setrun(p2, f'{idx:02d} / {total:02d}', 11, MUTED, bold=True)
    # marca
    tf3 = textbox(slide, 0.55, 6.74, 9.0, 0.3, MSO_ANCHOR.MIDDLE)
    setrun(tf3.paragraphs[0], deck, 9, RGBColor(0x9A,0xA6,0xAB))

def title_bar(slide, kicker, title):
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 1.18, fill=TEAL)
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 1.18, 13.333, 0.07, fill=CORAL)
    tf = textbox(slide, 0.55, 0.12, 12.2, 1.0, MSO_ANCHOR.MIDDLE)
    if kicker:
        pk = tf.paragraphs[0]
        setrun(pk, kicker.upper(), 12, CREAM, bold=True)
        pt = tf.add_paragraph()
    else:
        pt = tf.paragraphs[0]
    setrun(pt, title, 27, WHITE, bold=True)

def bg(slide):
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=WHITE)

def content_slide(prs, s, idx, total, deck, vbox=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title_bar(slide, s.get('k',''), s['t'])
    has_img = bool(s.get('img'))
    # bullets
    tw = 5.9 if has_img else 9.2
    tf = textbox(slide, 0.6, 1.55, tw, 5.0, MSO_ANCHOR.TOP)
    first = True
    for b in s['b']:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10); p.line_spacing = 1.06
        setrun(p, '●  ', 13, CORAL, bold=True)
        # soporta **negrita** simple al inicio "Titulo: resto"
        if '::' in b:
            head, rest = b.split('::', 1)
            setrun(p, head.strip()+': ', 16, TEAL, bold=True)
            setrun(p, rest.strip(), 16, INK)
        else:
            setrun(p, b, 16, INK)
    if has_img:
        # marco suave para la imagen
        add_image_fit(slide, s['img'], 6.75, 1.55, 6.1, 3.55)
    if vbox:
        video_box(slide)
    footer(slide, s['p'], idx, total, deck)
    # guion en notas
    if s.get('g'):
        slide.notes_slide.notes_text_frame.text = f"[{s['p']} · Slide {idx}] {s['t']}\n\n{s['g']}"
    return slide

def title_slide(prs, total, deck, sub, idx=1, presenter='Elmer'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=TEAL)
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 5.0, 13.333, 0.09, fill=CORAL)
    tf = textbox(slide, 0.8, 1.2, 11.7, 2.6, MSO_ANCHOR.BOTTOM)
    setrun(tf.paragraphs[0], 'APEX TWIN', 16, CREAM, bold=True)
    p = tf.add_paragraph(); setrun(p, 'SmartPark', 60, WHITE, bold=True)
    p2 = tf.add_paragraph(); setrun(p2, sub, 20, CREAM)
    # equipo
    tf2 = textbox(slide, 0.85, 5.25, 9.0, 1.9, MSO_ANCHOR.TOP)
    setrun(tf2.paragraphs[0], 'Equipo Apex Twin', 13, CREAM, bold=True)
    for key,(fn,role) in MEMBERS.items():
        pp = tf2.add_paragraph(); pp.space_after = Pt(1)
        setrun(pp, f'{fn}', 12, WHITE, bold=True)
        setrun(pp, f'  —  {role}', 12, CREAM)
    # curso
    tf3 = textbox(slide, 0.85, 0.5, 11.0, 0.6, MSO_ANCHOR.TOP)
    setrun(tf3.paragraphs[0], '1ASI0728 · Arquitecturas de Software Emergentes · Sección 10042', 12, CREAM)
    tf4 = textbox(slide, 9.7, 0.5, 3.0, 0.5, MSO_ANCHOR.TOP)
    pr = tf4.paragraphs[0]; pr.alignment = PP_ALIGN.RIGHT
    setrun(pr, deck, 12, CREAM, bold=True)
    video_box(slide)
    slide.notes_slide.notes_text_frame.text = f"[{presenter} · Slide {idx}] Portada\n\n(ver guion)"
    return slide

# =================================================================
#  KEYNOTE TB2
# =================================================================
DECK1 = 'SmartPark · Apex Twin — Keynote TB2'
A = ip('chapter-07'); C6 = ip('chapter-06'); C4 = ip('chapter-04')
kn = [
 dict(p='Elmer', k='El problema', t='Operar un estacionamiento a ciegas', b=[
   'Visibilidad:: el operador no conoce en tiempo real la ocupación ni el estado de seguridad del lote.',
   'Conductores:: pierden tiempo buscando plaza; congestión y mala experiencia.',
   'Seguridad:: los incidentes de humo se detectan tarde y sin contexto espacial.',
   'Datos dispersos:: sin un gemelo digital que unifique la operación física y digital.'], img=None,
   g='El problema que abordamos es claro: hoy un operador de estacionamiento trabaja casi a ciegas. No tiene una vista en tiempo real de cuántas plazas están ocupadas ni del estado de seguridad del lote. El conductor, por su lado, pierde tiempo buscando dónde estacionar. Y cuando ocurre un incidente, como humo o fuego, se detecta tarde y sin saber en qué zona exacta. La información está dispersa y no existe un gemelo digital que conecte el mundo físico con el digital.'),
 dict(p='Elmer', k='La solución', t='SmartPark: gemelo digital operativo', b=[
   'Cuatro productos digitales integrados:: Landing, Web App del operador, App móvil del conductor y el Gemelo digital + simulador IoT.',
   'Backend ASP.NET Core 8 con DDD, JWT y SignalR en Azure.',
   'Gemelo en Azure Digital Twins con visor 3D y ocupación en vivo.',
   'Cadena de seguridad:: humo geolocalizado y resaltado en el modelo 3D.'], img=os.path.join(A,'deployment-diagram.png'),
   g='Nuestra solución es SmartPark: un gemelo digital operativo del estacionamiento. La construimos como cuatro productos digitales integrados: la Landing Page, la aplicación web del operador, la app móvil del conductor y el gemelo digital alimentado por un simulador IoT. Por debajo, un backend en ASP.NET Core 8 con Domain-Driven Design, autenticación JWT y tiempo real con SignalR, todo desplegado en Azure. El corazón es el gemelo en Azure Digital Twins, que el operador ve en un visor 3D con la ocupación coloreada en vivo y las alertas de humo geolocalizadas.'),
 dict(p='Camila', k='Diseño', t='Design System y experiencia UX/UI', b=[
   'Sistema de diseño propio:: paleta teal/crema, tipografía Inter y componentes consistentes.',
   'Wireframes para los perfiles de operador y conductor.',
   'Guías de estilo aplicadas a web y móvil.',
   'Coherencia visual entre los cuatro productos.'], img=os.path.join(C6,'colores.png'),
   g='Desde el diseño, definimos un sistema propio para SmartPark: una paleta basada en teal y crema, la tipografía Inter y un conjunto de componentes reutilizables. Sobre esa base elaboramos los wireframes de los dos perfiles, el operador y el conductor, y las guías de estilo para web y móvil. El objetivo fue lograr coherencia visual entre los cuatro productos, de modo que la marca y la experiencia se sientan como un solo sistema.'),
 dict(p='Camila', k='Landing', t='Landing Page e Information Architecture', b=[
   'Propuesta de valor y secciones para operador y conductor.',
   'Planes de suscripción y navegación responsiva.',
   'Internacionalización inglés/español.',
   'Arquitectura de información:: labeling, navegación y SEO.'], img=os.path.join(A,'sprint-1-landing.png'),
   g='La Landing Page es la puerta de entrada. Comunica la propuesta de valor y tiene secciones específicas para operadores y conductores, además de los planes de suscripción y una navegación totalmente responsiva, con soporte de idioma inglés y español. Detrás trabajamos la arquitectura de información: el sistema de etiquetado, la navegación y las etiquetas SEO, para que el sitio sea claro tanto para las personas como para los buscadores.'),
 dict(p='Camila', k='Mockups', t='Del wireframe al producto', b=[
   'Mock-ups de alta fidelidad del dashboard y del gemelo digital.',
   'Pantallas de seguridad e incidentes y de eficiencia.',
   'Flujos de usuario del operador y del conductor.',
   'Validados contra el sistema de diseño.'], img=os.path.join(C6,'mockup-gemelo digital.png'),
   g='A partir de los wireframes elaboramos los mock-ups de alta fidelidad: el dashboard del operador, el gemelo digital, las pantallas de seguridad e incidentes y las de eficiencia. También diagramamos los flujos de usuario del operador y del conductor. Cada mock-up se validó contra el sistema de diseño, asegurando que lo que diseñamos fuera consistente con lo que finalmente se implementó.'),
 dict(p='Britney', k='Sprint 1', t='Planificación del Sprint 1', b=[
   'Plan de dos Sprints para cubrir el Product Backlog.',
   'Objetivo:: habilitar el flujo núcleo demostrable de extremo a extremo.',
   '94 Story Points comprometidos; User y Technical Stories priorizadas.',
   'Work-items asignados por frente de desarrollo.'], img=None,
   g='Organizamos el trabajo en dos Sprints. El objetivo del Sprint 1 fue habilitar el flujo núcleo de SmartPark de extremo a extremo: que una alerta de humo recorra toda la cadena, del simulador al gemelo y al dashboard, y que el conductor consulte disponibilidad desde la app. Comprometimos noventa y cuatro Story Points, seleccionando del Product Backlog las historias de usuario de mayor prioridad junto con las historias técnicas que las habilitan, y repartimos los work-items según el frente de cada integrante.'),
 dict(p='Britney', k='Gestión ágil', t='Tablero Trello y colaboración', b=[
   '29 work-items en tres listas: Done, To-Review e In-Process.',
   'Etiqueta de color por integrante; 25 tareas cerradas.',
   'Flujo feature-branches → develop → main (Pull Requests).',
   '200 commits en cinco repositorios de producto.'], img=os.path.join(A,'sprint-1-trello-board.png'),
   g='Gestionamos el Sprint en Trello, con tres listas: Done, To-Review e In-Process, y una etiqueta de color por integrante. Cerramos veinticinco de los veintinueve work-items. A nivel técnico seguimos un flujo de ramas ordenado: cada funcionalidad en su feature-branch, integrada a develop mediante Pull Request y liberada a main. En total registramos doscientos commits en los cinco repositorios de producto, lo que evidencia un trabajo colaborativo y trazable.'),
 dict(p='Britney', k='Móvil', t='App móvil del conductor', b=[
   'Construida en Power Apps (Canvas).',
   'Registro e inicio de sesión del conductor.',
   'Mapa de disponibilidad por zona y nivel.',
   'Registro de la ubicación del vehículo y conector custom desde el OpenAPI.'], img=os.path.join(C6,'Mob mock up - Home.png'),
   g='La experiencia del conductor vive en una app móvil construida en Power Apps. Permite registrarse e iniciar sesión, consultar un mapa de disponibilidad por zona y nivel, y registrar con un toque la ubicación del vehículo. Para conectarse al backend generamos un conector personalizado a partir del contrato OpenAPI del API, de modo que la app consume exactamente los mismos servicios que la web.'),
 dict(p='Abel', k='Gemelo 3D', t='Console del operador con gemelo en vivo', b=[
   'El gemelo digital 3D es el elemento central del console.',
   'Cada plaza se colorea según su ocupación: libre, ocupada o reservada.',
   'Panel de ocupación por zona con nivel de congestión.',
   'Auto-refresco y estados resilientes ante fallos.'], img=os.path.join(A,'sprint-1-console.png'),
   g='El console del operador pone al gemelo digital tridimensional en el centro. Cada plaza del modelo se colorea en vivo según su estado: verde si está libre, rojo si está ocupada y ámbar si está reservada, con datos reales que consulta al API. A la derecha, un panel resume la ocupación por zona y su nivel de congestión. Todo se auto-refresca y, si el gemelo no responde, la interfaz degrada con elegancia mostrando el último estado conocido.'),
 dict(p='Abel', k='Seguridad', t='Alertas de humo en tiempo real', b=[
   'Detección de humo geolocalizada en el modelo 3D.',
   'Resaltado de la zona afectada por el incidente.',
   'Push en vivo al dashboard mediante SignalR.',
   'Vista de simulador para las entrevistas de validación.'], img=os.path.join(A,'sprint-1-simulator.png'),
   g='La cadena de seguridad es uno de los diferenciales. Cuando el simulador emite una lectura de humo por encima del umbral, el incidente se geolocaliza y la zona afectada se resalta en el modelo 3D. La alerta llega al dashboard del operador en tiempo real mediante SignalR, sin necesidad de recargar. Para las entrevistas de validación implementamos además una vista de simulador, que permite disparar estos eventos en vivo y mostrar el comportamiento de extremo a extremo.'),
 dict(p='Abel', k='Web App', t='Aplicación web desplegada', b=[
   'Angular 20 con componentes standalone y signals.',
   'Login con JWT y sesión persistente.',
   'Ocupación por nivel con indicadores de congestión.',
   'Desplegada y verificada sobre Azure.'], img=os.path.join(A,'sprint-1-zones.png'),
   g='La aplicación web del operador está construida en Angular 20, con componentes standalone y signals para un estado reactivo y eficiente. Incluye inicio de sesión con JWT y sesión persistente, y vistas como la ocupación por nivel con indicadores de congestión. Lo importante es que no se quedó en local: está desplegada sobre Azure y verificamos el flujo completo de inicio de sesión y consulta de datos directamente sobre el entorno en la nube.'),
 dict(p='Hernán', k='Gemelo · IoT', t='Simulador IoT y Azure Digital Twins', b=[
   'Simulador Node.js de telemetría: ocupación y humo.',
   'Sincronización al grafo de twins vía JSON Patch.',
   'Capa anti-corrupción (ACL) sobre el SDK de ADT.',
   'Modelo 3D GLB con material único por plaza.'], img=os.path.join(A,'sprint-1-twin-occupancy.png'),
   g='Del lado del gemelo y el IoT, desarrollamos un simulador en Node.js que genera telemetría sintética de ocupación y de humo. Esa telemetría se sincroniza con el grafo de twins en Azure Digital Twins mediante operaciones JSON Patch, a través de una capa anti-corrupción que aísla nuestro dominio del SDK de Azure. Además generamos el modelo 3D en formato GLB con un material único por plaza, lo que permite colorear cada espacio individualmente según su ocupación en tiempo de ejecución.'),
 dict(p='Hernán', k='Despliegue', t='Desplegado en Azure', b=[
   'Landing y Web App como Storage static websites.',
   'API en App Service (.NET 8, Linux) con base SQLite embebida.',
   'Azure Digital Twins y Blob Storage para el modelo 3D.',
   'Despliegue con Az PowerShell; API documentada con Swagger.'], img=os.path.join(A,'sprint-1-deployment.png'),
   g='Cumplimos el hito de tener la primera versión desplegada. La Landing y la Web App se publican como sitios estáticos en Azure Storage; el API corre en App Service con .NET 8 sobre Linux, con una base SQLite embebida como solución funcional ante las restricciones de la suscripción de estudiante. El gemelo vive en Azure Digital Twins y el modelo 3D en Blob Storage. Todo el despliegue se ejecuta con Az PowerShell, que atraviesa el proxy de la universidad, y el API queda documentado con Swagger.'),
 dict(p='Hernán', k='Calidad y cierre', t='Validación, calidad y próximos pasos', b=[
   'Diseño de entrevistas de validación con ambos segmentos.',
   'Suite de pruebas: 41 métodos y 65 casos, en verde, con CI.',
   'Sprint 2:: flujo vehicular, eficiencia energética, sesión y costos, push.',
   '¡Gracias! SmartPark, el gemelo digital que opera por ti.'], img=os.path.join(A,'sprint-1-swagger.png'),
   g='Para cerrar: preparamos el diseño de las entrevistas de validación con ambos segmentos objetivo, y respaldamos la calidad con una suite de pruebas automatizadas de cuarenta y un métodos y sesenta y cinco casos, todos en verde y ejecutados en integración continua. El Sprint 2 abordará las funcionalidades complementarias: el flujo vehicular, la eficiencia energética, la gestión de sesión y costos, y la notificación push. Con esto demostramos el flujo núcleo de SmartPark, de extremo a extremo y sobre la nube. Muchas gracias.'),
]

def build_keynote():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    total = len(kn) + 1
    title_slide(prs, total, 'KEYNOTE TB2', 'Gemelo digital para la operación inteligente de estacionamientos', idx=1, presenter='Elmer')
    for i, s in enumerate(kn, start=2):
        content_slide(prs, s, i, total, DECK1, vbox=True)
    out = os.path.join(BASE, 'upc-pre-202601-1ASI0728-10042-ApexTwin-keynote-tb2.pptx')
    prs.save(out); return out, total

# =================================================================
#  ABOUT THE TEAM (v1)
# =================================================================
DECK2 = 'About the Team · Apex Twin — Sprint 1'

def at_title(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=TEAL)
    rect(slide, MSO_SHAPE.RECTANGLE, 0, 4.7, 13.333, 0.09, fill=CORAL)
    tf = textbox(slide, 0.9, 1.6, 11.5, 3.0, MSO_ANCHOR.BOTTOM)
    setrun(tf.paragraphs[0], 'APEX TWIN', 16, CREAM, bold=True)
    p = tf.add_paragraph(); setrun(p, 'About the Team', 52, WHITE, bold=True)
    p2 = tf.add_paragraph(); setrun(p2, 'Cómo construimos SmartPark — Sprint 1', 20, CREAM)
    tf2 = textbox(slide, 0.95, 5.0, 11.0, 1.8, MSO_ANCHOR.TOP)
    setrun(tf2.paragraphs[0], '1ASI0728 · Arquitecturas de Software Emergentes · Sección 10042', 12, CREAM)
    p3 = tf2.add_paragraph(); setrun(p3, 'Video About-the-Team — primera versión', 12, CREAM, bold=True)
    slide.notes_slide.notes_text_frame.text = 'Portada (voz en off · Elmer)'
    return slide

AT = [
 dict(p='Elmer', k='El equipo', t='Conoce a Apex Twin', kind='team', img=None,
   g='Somos Apex Twin, un equipo de cinco estudiantes del curso de Arquitecturas de Software Emergentes. Nos organizamos por frentes complementarios: liderazgo y backend, diseño y landing, aplicación móvil y validación, frontend del operador, y el gemelo digital con IoT. Cada integrante asumió la responsabilidad de su producto, pero trabajamos como una sola solución integrada.'),
 dict(p='Elmer', k='Proceso', t='Nuestra forma de trabajar', b=[
   'Marco ágil Scrum con una planificación de dos Sprints.',
   'Ramas feature → develop → main mediante Pull Requests.',
   'Coordinación diaria por Discord y tablero en Trello.',
   'Contrato de API (OpenAPI) acordado temprano para avanzar en paralelo.'], img=os.path.join(A,'sprint-1-trello-board.png'),
   g='Trabajamos con un marco ágil Scrum, planificando dos Sprints para cubrir todo el backlog. Adoptamos un flujo de ramas ordenado: cada funcionalidad en su rama, integrada a develop por Pull Request y liberada a main. Nos coordinamos a diario por Discord y gestionamos el avance en Trello. Una decisión clave fue acordar temprano el contrato del API con OpenAPI: eso nos permitió que el frontend web, la app móvil y el simulador avanzaran en paralelo sin bloquearse.'),
 dict(p='Elmer', k='Resultados', t='Lo que logramos en el Sprint 1', b=[
   'Cuatro productos digitales desplegados y verificados en Azure.',
   'Flujo núcleo de extremo a extremo: del simulador al gemelo 3D y al dashboard.',
   '200 commits y una suite de pruebas en verde.',
   'Outcome ABET: comunicación efectiva y trabajo colaborativo.'], img=os.path.join(A,'sprint-1-deployment.png'),
   g='Los resultados del Sprint 1 hablan por sí solos: cuatro productos digitales desplegados y verificados en Azure, con el flujo núcleo funcionando de extremo a extremo, desde el simulador IoT hasta el gemelo 3D y el dashboard del operador. Registramos doscientos commits y una suite de pruebas automatizadas en verde. Más allá de lo técnico, este proceso evidencia el outcome ABET de comunicación efectiva y de trabajo en un ambiente ágil y colaborativo.'),
 # testimonios
 dict(p='Elmer', k='Testimonio', t='Elmer — Team Leader & Backend', testimony=True, img=os.path.join(A,'sprint-1-swagger.png'),
   h=['Scrum Master: planificación y seguimiento de los Sprints.',
      'Backend ASP.NET Core 8 con DDD, JWT y SignalR.',
      'Persistencia EF Core y despliegue completo en Azure.'],
   g='Hola, soy Elmer Riva, líder del equipo. Como Scrum Master coordiné la planificación y el seguimiento de los Sprints. En lo técnico, construí el backend en ASP.NET Core 8 con Domain-Driven Design: la autenticación con JWT, los servicios de ocupación, la cadena de alertas con SignalR y la persistencia con Entity Framework. También me encargué del despliegue completo en Azure. Aprendí a diseñar una arquitectura multicomponente y a comunicar decisiones técnicas al equipo.'),
 dict(p='Camila', k='Testimonio', t='Camila — Landing & UI/UX', testimony=True, img=os.path.join(A,'sprint-1-landing.png'),
   h=['Sistema de diseño, wireframes y mock-ups.',
      'Arquitectura de información del sitio.',
      'Landing Page: propuesta de valor, planes e i18n.'],
   g='Soy Camila Sánchez y lideré el diseño y la Landing Page. Definí el sistema de diseño de SmartPark, los wireframes y los mock-ups de los dos perfiles, y trabajé la arquitectura de información del sitio. Implementé la landing con su propuesta de valor, los planes y el soporte de idiomas. Desarrollé mi competencia de diseño centrado en el usuario y de comunicar una identidad visual coherente entre productos.'),
 dict(p='Britney', k='Testimonio', t='Britney — App Móvil & Validación', testimony=True, img=os.path.join(C6,'Mob mock up - Home.png'),
   h=['App móvil del conductor en Power Apps.',
      'Conector custom desde el contrato OpenAPI.',
      'Apoyo en el diseño de entrevistas de validación.'],
   g='Hola, soy Britney Qqueso. Me encargué de la aplicación móvil del conductor en Power Apps: el registro, el inicio de sesión, el mapa de disponibilidad y el registro de la ubicación del vehículo, conectados al API por un conector personalizado. También apoyé el diseño de las entrevistas de validación. Aprendí a construir sobre plataformas low-code y a planificar la validación con usuarios reales.'),
 dict(p='Abel', k='Testimonio', t='Abel — Frontend del Operador', testimony=True, img=os.path.join(A,'sprint-1-console.png'),
   h=['Web del operador en Angular 20 (standalone, signals).',
      'Visor 3D del gemelo con ocupación en vivo.',
      'Tiempo real (SignalR) y alertas geolocalizadas.'],
   g='Soy Abel Valle y desarrollé la aplicación web del operador en Angular 20. Integré el visor 3D del gemelo digital con la ocupación coloreada en vivo, el dashboard con actualización en tiempo real por SignalR y las alertas de humo geolocalizadas. Cuidé la accesibilidad y los estados resilientes. Fortalecí mi competencia en frontend moderno y en integrar visualización 3D con datos en tiempo real.'),
 dict(p='Hernán', k='Testimonio', t='Hernán — Gemelo Digital & IoT', testimony=True, img=os.path.join(A,'sprint-1-twin-occupancy.png'),
   h=['Gemelo digital en Azure Digital Twins.',
      'Simulador IoT (ocupación y humo) en Node.js.',
      'Modelo 3D GLB con material único por plaza.'],
   g='Hola, soy Hernán Morales. Construí el gemelo digital en Azure Digital Twins y el simulador IoT en Node.js que lo alimenta con telemetría de ocupación y humo. Implementé la sincronización al grafo, la capa anti-corrupción y el modelo 3D con material por plaza. Aprendí a trabajar con tecnologías emergentes de gemelos digitales y a integrar IoT con la nube de forma desacoplada.'),
 dict(p='Elmer', k='Retrospectiva', t='Retrospectiva y próximos pasos', b=[
   'Aciertos:: contrato de API temprano y una única fuente de verdad del gemelo.',
   'Mejoras:: automatizar pruebas de integración y estandarizar el despliegue.',
   'Sprint 2:: flujo vehicular, eficiencia energética, sesión y costos, push.',
   '¡Gracias por acompañarnos en el proceso!'], img=os.path.join(A,'sprint-1-insights-web-services.png'),
   g='Para cerrar nuestra retrospectiva: como aciertos, acordar el contrato del API desde el inicio y mantener una única fuente de verdad para el gemelo nos permitió avanzar en paralelo sin choques. Como oportunidad de mejora, identificamos automatizar más pruebas de integración y estandarizar el despliegue. En el Sprint 2 abordaremos el flujo vehicular, la eficiencia energética, la gestión de sesión y costos, y la notificación push. Gracias por acompañarnos en el proceso de construir SmartPark.'),
]

def at_content(prs, s, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    title_bar(slide, s.get('k',''), s['t'])
    if s.get('kind') == 'team':
        tf = textbox(slide, 0.7, 1.6, 12.0, 5.0)
        first=True
        for key,(fn,role) in MEMBERS.items():
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            p.space_after = Pt(12)
            setrun(p, '●  ', 15, CORAL, bold=True)
            setrun(p, fn, 18, TEAL, bold=True)
            setrun(p, f'   —  {role}', 16, INK)
        video_box(slide)
    elif s.get('testimony'):
        add_image_fit(slide, s['img'], 0.7, 1.55, 5.7, 4.9)
        tf = textbox(slide, 6.7, 1.65, 6.2, 3.5)
        setrun(tf.paragraphs[0], 'Actividades en el Sprint 1', 14, CORAL, bold=True)
        for hh in s.get('h', []):
            p = tf.add_paragraph(); p.space_before=Pt(8); p.line_spacing=1.06
            setrun(p, '●  ', 13, CORAL, bold=True)
            setrun(p, hh, 15, INK)
        tf2 = textbox(slide, 6.7, 4.95, 3.3, 0.5)
        setrun(tf2.paragraphs[0], 'Testimonio completo: ver guion (notas / TXT).', 11, MUTED)
        video_box(slide)
    else:
        has_img = bool(s.get('img'))
        tw = 5.9 if has_img else 9.2
        tf = textbox(slide, 0.6, 1.6, tw, 4.9)
        first=True
        for b in s['b']:
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            p.space_after = Pt(10); p.line_spacing=1.06
            setrun(p, '●  ', 13, CORAL, bold=True)
            if '::' in b:
                head, rest = b.split('::',1)
                setrun(p, head.strip()+': ', 16, TEAL, bold=True); setrun(p, rest.strip(), 16, INK)
            else:
                setrun(p, b, 16, INK)
        if has_img:
            add_image_fit(slide, s['img'], 6.75, 1.55, 6.1, 3.7)
        if s.get('testimony') or s.get('kind')=='team':
            video_box(slide)
    footer(slide, s['p'], idx, total, DECK2)
    slide.notes_slide.notes_text_frame.text = f"[{s['p']} · Slide {idx}] {s['t']}\n\n{s['g']}"
    return slide

def build_about_team():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    total = len(AT) + 1
    at_title(prs, total)
    for i, s in enumerate(AT, start=2):
        at_content(prs, s, i, total)
    out = os.path.join(BASE, 'upc-pre-202601-1ASI0728-10042-ApexTwin-about-the-team-sprint-1.pptx')
    prs.save(out); return out, total

# =================================================================
#  GUIONES (.txt)
# =================================================================
def write_guion_keynote(total):
    lines = []
    lines.append('GUION — KEYNOTE TB2 · SmartPark (Apex Twin)')
    lines.append('1ASI0728 Arquitecturas de Software Emergentes · Sección 10042')
    lines.append('='*70)
    lines.append('')
    lines.append('Reparto de slides (números consecutivos por integrante):')
    lines.append('  Elmer   -> 1, 2, 3      Camila  -> 4, 5, 6')
    lines.append('  Britney -> 7, 8, 9      Abel    -> 10, 11, 12')
    lines.append('  Hernán  -> 13, 14, 15')
    lines.append('Nota: en cada slide hay un recuadro inferior derecho para el video del expositor.')
    lines.append('='*70)
    lines.append('')
    fn = MEMBERS['Elmer'][0]
    lines.append(f'--- SLIDE 1 · {fn} · PORTADA ---')
    lines.append('Buenas, somos el equipo Apex Twin y les presentamos SmartPark, un gemelo digital '
                 'para la operación inteligente de estacionamientos. Esta es la entrega TB2, donde '
                 'mostramos la primera versión desplegada de nuestra solución y la ejecución del Sprint 1. '
                 'Mi nombre es Elmer Riva y soy el líder del equipo; a lo largo de la presentación '
                 'cada integrante expondrá su parte.')
    lines.append('')
    for i, s in enumerate(kn, start=2):
        fn = MEMBERS[s['p']][0]
        lines.append(f"--- SLIDE {i} · {fn} · {s['t']} ---")
        lines.append(s['g'])
        lines.append('')
    open(os.path.join(BASE,'upc-pre-202601-1ASI0728-10042-ApexTwin-keynote-tb2-guion.txt'),'w',encoding='utf-8').write('\n'.join(lines))

def write_guion_about(total):
    lines = []
    lines.append('GUION — VIDEO ABOUT-THE-TEAM (v1) · SmartPark (Apex Twin)')
    lines.append('1ASI0728 Arquitecturas de Software Emergentes · Sección 10042')
    lines.append('='*70)
    lines.append('Estructura: retrospectiva del proceso (voz en off) + 1 testimonio por integrante.')
    lines.append('Sugerencia de timing: ~5 min de proceso/retrospectiva + ~1 min por testimonio.')
    lines.append('='*70)
    lines.append('')
    lines.append('### NARRACIÓN / VOZ EN OFF (proceso del equipo) — sugerido: Elmer (líder)')
    lines.append('')
    lines.append('[Slide 1 · Portada]')
    lines.append('Somos Apex Twin y este es el video About-the-Team: cómo construimos SmartPark durante el Sprint 1.')
    lines.append('')
    for s in AT[:3]:
        lines.append(f"[{s['t']}]")
        lines.append(s['g']); lines.append('')
    lines.append('='*70)
    lines.append('### TESTIMONIOS POR INTEGRANTE (ante cámara, ~1 min c/u)')
    lines.append('')
    order = ['Elmer','Camila','Britney','Abel','Hernán']
    tmap = {s['p']: s for s in AT if s.get('testimony')}
    for key in order:
        fn, role = MEMBERS[key]
        s = tmap[key]
        lines.append(f'--- {fn} · {role} ---')
        lines.append(s['g']); lines.append('')
    lines.append('='*70)
    lines.append('### CIERRE / RETROSPECTIVA — voz en off (Elmer)')
    lines.append('')
    s = AT[-1]
    lines.append(s['g']); lines.append('')
    open(os.path.join(BASE,'upc-pre-202601-1ASI0728-10042-ApexTwin-about-the-team-sprint-1-guion.txt'),'w',encoding='utf-8').write('\n'.join(lines))

if __name__ == '__main__':
    o1, t1 = build_keynote();       print('KEYNOTE   ->', o1, '(', t1, 'slides )')
    o2, t2 = build_about_team();    print('ABOUT-TEAM->', o2, '(', t2, 'slides )')
    write_guion_keynote(t1);        print('guion keynote OK')
    write_guion_about(t2);          print('guion about-team OK')
